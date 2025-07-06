# app.py – EduScribe AI: Lecture Notes Generator from Slides & Audio

import streamlit as st
from pptx import Presentation
import whisper
import ollama
import os

# --- Slide Extraction ---
def extract_text_from_pptx(uploaded_pptx):
    prs = Presentation(uploaded_pptx)
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slide_texts.append((f"Slide {i+1}", text))
    return slide_texts

# --- Audio Transcription ---
def transcribe_audio(file):
    model = whisper.load_model("base")  # You can use 'tiny', 'base', 'small' etc.
    audio_path = f"temp_audio.{file.type.split('/')[-1]}"
    with open(audio_path, "wb") as f:
        f.write(file.read())
    result = model.transcribe(audio_path)
    os.remove(audio_path)
    return result["text"]

# --- Generate LLM Prompt ---
def build_prompt(slides, transcript):
    combined_text = "\n\n".join([f"{title}: {content}" for title, content in slides])
    prompt = f"""
You are an education assistant. Given the following lecture slides and audio transcript, generate:
1. Summarized lecture notes
2. 3-5 Q&A pairs for revision
3. Key terms with short definitions

Slides:
{combined_text[:2000]}

Transcript:
{transcript[:2000]}

Respond in Markdown format.
"""
    return prompt

# --- Query LLM (Ollama) ---
def query_llm(prompt):
    response = ollama.chat(model="llama3", messages=[{"role": "user", "content": prompt}])
    return response['message']['content']

# --- Streamlit UI ---
st.set_page_config(page_title="EduScribe AI", layout="wide")
st.title("🎓 EduScribe AI – Lecture Notes Generator")

pptx_file = st.file_uploader("Upload Lecture Slides (PPTX)", type=["pptx"])
audio_file = st.file_uploader("Upload Lecture Audio (MP3/WAV)", type=["mp3", "wav"])

if pptx_file and audio_file:
    with st.spinner("Processing slides and audio..."):
        slide_data = extract_text_from_pptx(pptx_file)
        audio_transcript = transcribe_audio(audio_file)
        st.success("✅ Slide content and audio transcription ready.")

    if st.button("🧠 Generate Notes"):
        with st.spinner("Generating notes with LLM..."):
            prompt = build_prompt(slide_data, audio_transcript)
            result = query_llm(prompt)
            st.markdown("### 📚 Generated Lecture Notes")
            st.markdown(result)
else:
    st.info("Upload both PPTX and audio files to begin.")
